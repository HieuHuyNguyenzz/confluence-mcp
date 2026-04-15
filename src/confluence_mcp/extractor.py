"""Extract content from attachment files and convert to Markdown."""

import csv
import io
import json
import zipfile
import tarfile
from typing import Any

from bs4 import BeautifulSoup
from markdownify import markdownify as md


class FileExtractor:
    """Extract content from attachment files and return as Markdown."""

    @staticmethod
    def extract(filename: str, file_bytes: bytes) -> str:
        """Extract content based on file extension.

        Returns extracted content as Markdown string.
        Returns error message if extraction fails.
        Returns unsupported message if format not recognized.
        """
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        extractors: dict[str, Any] = {
            "pdf": _extract_pdf,
            "docx": _extract_docx,
            "doc": _extract_doc,
            "xlsx": _extract_xlsx,
            "xls": _extract_xls,
            "csv": _extract_csv,
            "pptx": _extract_pptx,
            "ppt": _extract_ppt,
            "txt": _extract_text,
            "md": _extract_text,
            "json": _extract_json,
            "xml": _extract_xml,
            "html": _extract_html,
            "rtf": _extract_text,
            "log": _extract_text,
            "yaml": _extract_text,
            "yml": _extract_text,
            "toml": _extract_text,
            "ini": _extract_text,
            "cfg": _extract_text,
            "sh": _extract_text,
            "py": _extract_text,
            "js": _extract_text,
            "ts": _extract_text,
            "java": _extract_text,
            "c": _extract_text,
            "cpp": _extract_text,
            "h": _extract_text,
            "go": _extract_text,
            "rs": _extract_text,
            "rb": _extract_text,
            "php": _extract_text,
            "sql": _extract_text,
            "css": _extract_text,
            "zip": _extract_zip,
            "rar": _extract_rar,
            "7z": _extract_7z,
            "tar": _extract_tar,
            "gz": _extract_tar,
            "xmind": _extract_xmind,
            "drawio": _extract_drawio,
            "msg": _extract_msg,
        }

        extractor = extractors.get(ext)
        if extractor:
            try:
                return extractor(file_bytes)
            except Exception as e:
                return f"[Error extracting {filename}: {type(e).__name__}: {e}]"

        return f"[Unsupported file type: .{ext} ({_guess_media_type(ext)})]"


def _extract_pdf(data: bytes) -> str:
    """Extract text from PDF files."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"## Page {i}\n\n{text.strip()}")
    return "\n\n".join(pages) if pages else "[No text content in PDF]"


def _extract_docx(data: bytes) -> str:
    """Extract content from .docx files as Markdown."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = []

    for para in doc.paragraphs:
        if not para.text.strip():
            continue

        style = para.style.name.lower() if para.style else ""

        if "heading" in style:
            try:
                level = int(style.replace("heading", "").replace(" ", "").strip())
            except ValueError:
                level = 1
            parts.append(f"{'#' * level} {para.text.strip()}")
        elif "list" in style:
            parts.append(f"- {para.text.strip()}")
        else:
            parts.append(para.text.strip())

    for table in doc.tables:
        md_table = _table_to_markdown(table)
        if md_table:
            parts.append(md_table)

    return "\n\n".join(p for p in parts if p) if parts else "[No text content]"


def _extract_doc(data: bytes) -> str:
    """Legacy .doc files - best effort text extraction."""
    try:
        text = data.decode("utf-8", errors="replace")
        return text.strip()
    except Exception:
        return "[Cannot extract .doc file - convert to .docx for better support]"


def _extract_xlsx(data: bytes) -> str:
    """Extract content from .xlsx files as Markdown tables."""
    from openpyxl import load_workbook

    wb = load_workbook(read_only=True, data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))

        if not rows:
            continue

        table_lines = []
        table_lines.append(f"## {sheet_name}")

        headers = [str(c) if c is not None else "" for c in rows[0]]
        table_lines.append("| " + " | ".join(headers) + " |")
        table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

        for row in rows[1:]:
            if all(c is None for c in row):
                continue
            cells = [str(c) if c is not None else "" for c in row]
            table_lines.append("| " + " | ".join(cells) + " |")

        sheets.append("\n".join(table_lines))

    return "\n\n".join(sheets) if sheets else "[Empty spreadsheet]"


def _extract_xls(data: bytes) -> str:
    """Legacy .xls files - best effort."""
    return "[Cannot extract .xls file - convert to .xlsx for better support]"


def _extract_csv(data: bytes) -> str:
    """Extract content from CSV files as Markdown tables."""
    text = data.decode("utf-8", errors="replace")

    sniffer = csv.Sniffer()
    try:
        dialect = sniffer.sniff(text[:4096])
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(text), dialect)
    rows = list(reader)

    if not rows:
        return "[Empty CSV]"

    table_lines = []
    headers = rows[0]
    table_lines.append("| " + " | ".join(headers) + " |")
    table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

    for row in rows[1:]:
        if not any(row):
            continue
        table_lines.append("| " + " | ".join(row) + " |")

    return "\n".join(table_lines)


def _extract_pptx(data: bytes) -> str:
    """Extract content from .pptx files as Markdown."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        parts = [f"## Slide {i}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())

            if shape.has_table:
                table_lines = []
                table = shape.table
                if table:
                    for row_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        if row_idx == 0:
                            table_lines.append("| " + " | ".join(cells) + " |")
                            table_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        else:
                            table_lines.append("| " + " | ".join(cells) + " |")

                if table_lines:
                    parts.append("\n".join(table_lines))

        if len(parts) > 1:
            slides.append("\n\n".join(parts))

    return "\n\n".join(slides) if slides else "[No text content]"


def _extract_ppt(data: bytes) -> str:
    """Legacy .ppt files."""
    return "[Cannot extract .ppt file - convert to .pptx for better support]"


def _extract_text(data: bytes) -> str:
    """Return plain text content."""
    return data.decode("utf-8", errors="replace").strip()


def _extract_json(data: bytes) -> str:
    """Parse and reformat JSON as readable text."""
    text = data.decode("utf-8", errors="replace")
    try:
        parsed = json.loads(text)
        return json.dumps(parsed, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return text.strip()


def _extract_xml(data: bytes) -> str:
    """Parse and format XML."""
    from xml.dom import minidom

    text = data.decode("utf-8", errors="replace")
    try:
        dom = minidom.parseString(text)
        return dom.toprettyxml(indent="  ")
    except Exception:
        return text.strip()


def _extract_html(data: bytes) -> str:
    """Convert HTML to Markdown."""
    text = data.decode("utf-8", errors="replace")
    soup = BeautifulSoup(text, "html.parser")
    return md(str(soup), heading_style="ATX", bullets="-").strip()


def _table_to_markdown(table) -> str:
    """Convert a python-docx table to Markdown."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _guess_media_type(ext: str) -> str:
    """Guess media type from extension."""
    types = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "svg": "image/svg+xml",
        "webp": "image/webp",
        "mp4": "video/mp4",
        "avi": "video/x-msvideo",
        "mp3": "audio/mpeg",
        "wav": "audio/wav",
        "zip": "application/zip",
        "rar": "application/x-rar-compressed",
        "7z": "application/x-7z-compressed",
        "exe": "application/x-executable",
        "dll": "application/x-dosexec",
        "iso": "application/x-iso9660-image",
    }
    return types.get(ext, "application/octet-stream")


def _extract_zip(data: bytes) -> str:
    """Extract content from ZIP archives."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            parts = []
            for name in z.namelist():
                if name.endswith("/"):
                    continue
                with z.open(name) as f:
                    content_bytes = f.read()
                    extracted = FileExtractor.extract(name, content_bytes)
                    parts.append(f"### File: {name}\n\n{extracted}")
            return "\n\n---\n\n".join(parts) if parts else "[ZIP archive is empty]"
    except Exception as e:
        return f"[Error extracting ZIP: {e}]"


def _extract_tar(data: bytes) -> str:
    """Extract content from TAR or TAR.GZ archives."""
    try:
        # tarfile.open handles .tar, .tar.gz, .tar.bz2, .tar.xz automatically
        with tarfile.open(fileobj=io.BytesIO(data), mode="r:*") as t:
            parts = []
            for member in t.getmembers():
                if member.isfile():
                    f = t.extractfile(member)
                    if f:
                        content_bytes = f.read()
                        extracted = FileExtractor.extract(member.name, content_bytes)
                        parts.append(f"### File: {member.name}\n\n{extracted}")
            return "\n\n---\n\n".join(parts) if parts else "[TAR archive is empty]"
    except Exception as e:
        return f"[Error extracting TAR: {e}]"


def _extract_rar(data: bytes) -> str:
    """Extract content from RAR archives."""
    try:
        import rarfile
        with rarfile.RarFile(io.BytesIO(data)) as r:
            parts = []
            for name in r.namelist():
                if name.endswith("/"):
                    continue
                with r.open(name) as f:
                    content_bytes = f.read()
                    extracted = FileExtractor.extract(name, content_bytes)
                    parts.append(f"### File: {name}\n\n{extracted}")
            return "\n\n---\n\n".join(parts) if parts else "[RAR archive is empty]"
    except ImportError:
        return "[Error: 'rarfile' library not installed. Please install it to extract RAR files.]"
    except Exception as e:
        return f"[Error extracting RAR: {e}]"


def _extract_7z(data: bytes) -> str:
    """Extract content from 7z archives."""
    try:
        import py7zr
        with py7zr.SevenZipFile(io.BytesIO(data), mode='r') as z:
            parts = []
            # py7zr.SevenZipFile.readall() returns a dict {filename: Bio}
            all_files = z.readall()
            for name, bio in all_files.items():
                content_bytes = bio.read()
                extracted = FileExtractor.extract(name, content_bytes)
                parts.append(f"### File: {name}\n\n{extracted}")
            return "\n\n---\n\n".join(parts) if parts else "[7z archive is empty]"
    except ImportError:
        return "[Error: 'py7zr' library not installed. Please install it to extract 7z files.]"
    except Exception as e:
        return f"[Error extracting 7z: {e}]"


def _extract_xmind(data: bytes) -> str:
    """Extract content from XMind files as Markdown."""
    try:
        from xmindparser import xmind_to_dict

        with open("/tmp/temp_xmind.xmind", "wb") as f:
            f.write(data)

        result = xmind_to_dict("/tmp/temp_xmind.xmind")
        
        import os
        os.remove("/tmp/temp_xmind.xmind")

        return _format_xmind_to_markdown(result)
    except ImportError:
        return "[Error: 'xmindparser' library not installed. Please install it to extract XMind files.]"
    except Exception as e:
        return f"[Error extracting XMind: {type(e).__name__}: {e}]"


def _format_xmind_to_markdown(xmind_data: dict) -> str:
    """Format XMind data to Markdown hierarchy."""
    parts = []

    def process_topic(topic: dict, level: int = 1) -> list:
        lines = []
        title = topic.get("title", "Untitled")
        
        heading = "#" * min(level, 6)
        lines.append(f"{heading} {title}")
        
        notes = topic.get("notes", {})
        if notes:
            note_content = notes.get("plain", "")
            if note_content:
                lines.append(f"\n_{note_content}_")
        
        labels = topic.get("labels", [])
        if labels:
            lines.append(f"\nLabels: {', '.join(labels)}")
        
        topics = topic.get("topics", [])
        if topics:
            for subtopic in topics:
                sub_lines = process_topic(subtopic, level + 1)
                lines.extend(sub_lines)
        
        return lines

    if isinstance(xmind_data, list):
        for sheet in xmind_data:
            sheet_title = sheet.get("title", "Sheet")
            parts.append(f"## {sheet_title}\n")
            
            root_topic = sheet.get("topic", {})
            if root_topic:
                content = process_topic(root_topic, 2)
                parts.extend(content)
                
            parts.append("")
    elif isinstance(xmind_data, dict):
        root_topic = xmind_data.get("topic", {})
        if root_topic:
            content = process_topic(root_topic, 1)
            parts.extend(content)

    return "\n".join(parts) if parts else "[Empty XMind file]"


def _extract_drawio(data: bytes) -> str:
    """Extract content from drawio XML files as Markdown."""
    try:
        text = data.decode("utf-8", errors="replace")
        soup = BeautifulSoup(text, "xml")
        
        parts = []
        parts.append("## Diagram")
        
        mxgraph = soup.find("mxfile")
        if mxgraph:
            diagram_title = mxgraph.get("name", "Untitled Diagram")
            parts.append(f"**Name:** {diagram_title}\n")
        
        for diagram in soup.find_all("diagram"):
            diagram_name = diagram.get("name", "")
            if diagram_name:
                parts.append(f"\n### {diagram_name}")
            
            diagram_xml = diagram.string or ""
            if diagram_xml:
                inner_soup = BeautifulSoup(diagram_xml, "xml")
                
                for cell in inner_soup.find_all("mxCell"):
                    if cell.get("vertex") == "1":
                        parent = cell.get("parent")
                        style = cell.get("style", "")
                        value = cell.get("value", "")
                        
                        if value and value != "":
                            value_clean = BeautifulSoup(value, "html.parser").get_text()
                            if "shape=" in style:
                                if "rhombus" in style:
                                    parts.append(f"- ◇ {value_clean}")
                                elif "ellipse" in style:
                                    parts.append(f"- ○ {value_clean}")
                                elif "rect" in style or "rounded" in style:
                                    parts.append(f"- □ {value_clean}")
                                else:
                                    parts.append(f"- {value_clean}")
                            else:
                                parts.append(f"- {value_clean}")
                
                for edge in inner_soup.find_all("mxCell"):
                    if edge.get("edge") == "1":
                        style = edge.get("style", "")
                        value = edge.get("value", "")
                        source = edge.get("source", "")
                        target = edge.get("target", "")
                        
                        if value or source or target:
                            label = value if value else "→"
                            parts.append(f"- **{source}** → **{target}**: {label}")
        
        return "\n".join(parts) if parts else "[Empty drawio diagram]"
    except Exception as e:
        return f"[Error extracting drawio: {type(e).__name__}: {e}]"


def _extract_msg(data: bytes) -> str:
    """Extract content from Outlook .msg files as Markdown."""
    try:
        from extract_msg import Message

        with open("/tmp/temp_msg.msg", "wb") as f:
            f.write(data)

        msg = Message("/tmp/temp_msg.msg")
        
        parts = []
        
        subject = msg.subject or "No Subject"
        parts.append(f"## Email: {subject}\n")
        
        parts.append("**From:**")
        sender = msg.sender or "Unknown"
        parts.append(f"- {sender}")
        
        to_recipients = msg.to or ""
        if to_recipients:
            parts.append(f"\n**To:**")
            for recipient in to_recipients.split(";"):
                parts.append(f"- {recipient.strip()}")
        
        cc_recipients = msg.cc
        if cc_recipients:
            parts.append(f"\n**CC:**")
            for recipient in cc_recipients.split(";"):
                parts.append(f"- {recipient.strip()}")
        
        date = msg.date
        if date:
            parts.append(f"\n**Date:** {date}")
        
        body = msg.body
        if body:
            parts.append(f"\n---\n\n### Body\n\n{body}")
        
        html_body = msg.htmlBody
        if html_body:
            try:
                html_soup = BeautifulSoup(html_body, "html.parser")
                text_body = html_soup.get_text(separator="\n", strip=True)
                if text_body and not body:
                    parts.append(f"\n---\n\n### Body (HTML)\n\n{text_body}")
            except:
                pass
        
        attachments = msg.attachments
        if attachments:
            parts.append(f"\n### Attachments\n")
            for att in attachments:
                att_name = att.name if hasattr(att, 'name') else str(att)
                parts.append(f"- {att_name}")
        
        msg.close()
        
        import os
        os.remove("/tmp/temp_msg.msg")
        
        return "\n".join(parts)
    except ImportError:
        return "[Error: 'extract-msg' library not installed. Please install it to extract MSG files.]"
    except Exception as e:
        return f"[Error extracting MSG: {type(e).__name__}: {e}]"
